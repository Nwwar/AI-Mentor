import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.embeddings import HuggingFaceInstructEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
from langchain_community.llms import HuggingFaceHub


def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text


def PPTX_Reader(pptx_docs):
    text = []
    for pptx in pptx_docs:
        prs = Presentation(pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

    return "\n".join(text)


def Word_Reader(docx_docs):
    full_text = []
    for docx in docx_docs:
        doc = Document(docx)
        for para in doc.paragraphs:
            full_text.append(para.text)
    return '\n'.join(full_text)


def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore(text_chunks):
    embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature": 0.5, "max_length": 512})

    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain


def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)


def main():
    load_dotenv()
    st.set_page_config(page_title="Chat with multiple documents",
                       page_icon=":books:")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.header("Chat with multiple documents :books:")
    user_question = st.text_input("Ask a question about your documents:")
    if user_question:
        handle_userinput(user_question)

    with st.sidebar:
        st.subheader("Your documents")
        uploaded_file = st.file_uploader(
            "Upload your PDFs here and click on 'Process'", accept_multiple_files=True, type=['pdf', 'pptx', 'docx'])
        if st.button("Process"):
            with st.spinner("Processing"):
                # get text
                if uploaded_file:
                    pdf_docs = [file for file in uploaded_file if file.name.endswith(
                        ".pdf")]
                    pptx_docs = [file for file in uploaded_file if file.name.endswith(
                        ".pptx")]
                    docx_docs = [file for file in uploaded_file if file.name.endswith(
                        ".docx")]
                    pdf_text = get_pdf_text(pdf_docs)
                    pptx_text = PPTX_Reader(pptx_docs)
                    docx_text = Word_Reader(docx_docs)
                    raw_text = pdf_text + pptx_text + docx_text

                # get the text chunks
                text_chunks = get_text_chunks(raw_text)

                # create vector store
                vectorstore = get_vectorstore(text_chunks)

                # create conversation chain
                st.session_state.conversation = get_conversation_chain(
                    vectorstore)


if __name__ == '__main__':
    main()
